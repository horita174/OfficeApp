from pptx import Presentation

# prs = Presentation()
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"


# prs.save('test.pptx')
prs = Presentation("test.pptx")
# スライドのレイアウトを選択する（0以外の数値を使用する）
title_slide_layout = prs.slide_layouts[1]

# 新しいスライドを追加する
slide = prs.slides.add_slide(title_slide_layout)

# タイトルとサブタイトルを変更する
title = slide.shapes.title
title.text = "新しいタイトル"

subtitle = slide.placeholders[1]

# テキストボックスを追加
left = top = width = height = 10  # テキストボックスの位置とサイズを指定
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame


# パワーポイントファイル内の文字数をカウント
character_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    character_count += len(run.text)


# # 文字数をテキストボックスに挿入
# slide = prs.slides[-1]  # 最後のスライドを取得
# textbox = slide.shapes.add_textbox(left, top + height, width, height)
# text_frame = textbox.text_frame
# text_frame.text = 
subtitle.text = "文字数: {}".format(character_count)

# 変更内容を保存する
prs.save("test2.pptx")